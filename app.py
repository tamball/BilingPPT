import io
import os
import re
import shutil
import subprocess
import tempfile

import openai
import streamlit as st
from deep_translator import GoogleTranslator
from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt, Emu

# ── Language options ──────────────────────────────────────────────────────────

LANGUAGES = {
    "English": "en",
    "Chinese (Simplified)": "zh-CN",
    "Chinese (Traditional)": "zh-TW",
    "Spanish": "es",
    "French": "fr",
    "Korean": "ko",
    "Japanese": "ja",
    "Portuguese": "pt",
    "Arabic": "ar",
    "Hindi": "hi",
    "German": "de",
    "Italian": "it",
    "Russian": "ru",
    "Vietnamese": "vi",
    "Thai": "th",
    "Indonesian": "id",
    "Malay": "ms",
    "Tagalog": "tl",
    "Ukrainian": "uk",
    "Dutch": "nl",
}

# ── Theme palette ─────────────────────────────────────────────────────────────

THEMES = {
    "Dark (Navy)": {
        "bg":      RGBColor(0x1A, 0x1A, 0x2E),
        "orig":    RGBColor(0xFF, 0xFF, 0xFF),
        "trans":   RGBColor(0xFF, 0xD7, 0x00),
        "divider": RGBColor(0x44, 0x44, 0x66),
        "label":   RGBColor(0x88, 0x88, 0xAA),
    },
    "Dark (Blue)": {
        "bg":      RGBColor(0x0D, 0x2B, 0x55),
        "orig":    RGBColor(0xFF, 0xFF, 0xFF),
        "trans":   RGBColor(0x7E, 0xC8, 0xE3),
        "divider": RGBColor(0x1E, 0x4D, 0x8C),
        "label":   RGBColor(0x7E, 0xC8, 0xE3),
    },
    "Light": {
        "bg":      RGBColor(0xF8, 0xF8, 0xF8),
        "orig":    RGBColor(0x1A, 0x1A, 0x2E),
        "trans":   RGBColor(0xB0, 0x30, 0x20),
        "divider": RGBColor(0xCC, 0xCC, 0xCC),
        "label":   RGBColor(0x88, 0x88, 0x88),
    },
    "Black": {
        "bg":      RGBColor(0x00, 0x00, 0x00),
        "orig":    RGBColor(0xFF, 0xFF, 0xFF),
        "trans":   RGBColor(0x00, 0xFF, 0xCC),
        "divider": RGBColor(0x33, 0x33, 0x33),
        "label":   RGBColor(0x66, 0x66, 0x66),
    },
}

# ── Text processing ───────────────────────────────────────────────────────────

# Languages whose scripts are full-width (each glyph ≈ 1 em wide).
# Latin scripts average ~0.55 em per character.
CJK_LANG_CODES = {"zh-CN", "zh-TW", "ja", "ko"}

def _char_width_factor(lang_code: str) -> float:
    """Return estimated average character width as a fraction of font size (em)."""
    return 1.0 if lang_code in CJK_LANG_CODES else 0.55

def _find_libreoffice() -> str | None:
    """Return the LibreOffice executable path, or None if not found."""
    candidates = [
        "libreoffice",
        "soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for cmd in candidates:
        if shutil.which(cmd) or os.path.isfile(cmd):
            return cmd
    return None


def extract_paragraphs(uploaded_file) -> list[str]:
    """Return non-empty paragraph strings from a .doc or .docx file."""
    filename = uploaded_file.name.lower()

    if filename.endswith(".docx"):
        doc = Document(uploaded_file)
        return [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    # ── .doc (legacy binary format) ──────────────────────────────────────────
    lo_cmd = _find_libreoffice()
    if lo_cmd is None:
        raise RuntimeError(
            "Cannot open .doc files: LibreOffice is not installed. "
            "Please install LibreOffice (https://www.libreoffice.org) "
            "or save your file as .docx and re-upload."
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        doc_path = os.path.join(tmpdir, "input.doc")
        with open(doc_path, "wb") as f:
            f.write(uploaded_file.read())

        result = subprocess.run(
            [lo_cmd, "--headless", "--convert-to", "docx", "--outdir", tmpdir, doc_path],
            capture_output=True,
            timeout=60,
        )

        docx_path = os.path.join(tmpdir, "input.docx")
        if not os.path.exists(docx_path):
            stderr = result.stderr.decode(errors="replace")
            raise RuntimeError(f"LibreOffice conversion failed: {stderr}")

        doc = Document(docx_path)
        return [p.text.strip() for p in doc.paragraphs if p.text.strip()]


def split_sentences(text: str) -> list[str]:
    """Split text into sentences on common terminal punctuation."""
    parts = re.split(r'(?<=[.!?。！？])\s+', text)
    return [s.strip() for s in parts if s.strip()]


def max_chars_for_font(font_pt: int, lang_code: str = "en") -> int:
    """
    Estimate the maximum characters that safely fit in one text box half,
    accounting for CJK full-width glyphs vs narrow Latin glyphs.
    """
    BOX_H_PT   = (7.5 - 0.4) / 2 * 72
    lines      = int(BOX_H_PT / (font_pt * 1.2))
    usable_w   = (13.33 - 0.9) * 72
    factor     = _char_width_factor(lang_code)
    chars_line = int(usable_w / (font_pt * factor))
    return max(10, int(lines * chars_line * 0.90))


def optimal_font_size(chunks: list[str], lang_code: str) -> int:
    """
    Return the largest font size (pt) at which the longest chunk still fits
    within one slide-half text box.
    """
    if not chunks:
        return 54
    max_len = max(len(c) for c in chunks)
    for f in range(80, 34, -1):
        if max_chars_for_font(f, lang_code) >= max_len:
            return f
    return 35


def _src_max_chars(orig_font_pt: int, tgt_font_pt: int, src_code: str, tgt_code: str) -> int:
    """
    Max source-text characters per slide, accounting for cross-script expansion
    so that the *translated* text also fits in its half of the slide.

    CJK → Latin translations expand ~2.5×; Latin → CJK compress to ~0.4×.
    Using the inverse expansion as a divisor keeps translated text in-bounds.
    """
    src_max = max_chars_for_font(orig_font_pt, src_code)
    tgt_max = max_chars_for_font(tgt_font_pt, tgt_code)

    src_is_cjk = src_code in CJK_LANG_CODES
    tgt_is_cjk = tgt_code in CJK_LANG_CODES

    if src_is_cjk and not tgt_is_cjk:
        # e.g. Chinese → English: translation expands ~2.5×, cap src accordingly
        return min(src_max, int(tgt_max / 2.5))
    elif not src_is_cjk and tgt_is_cjk:
        # e.g. English → Chinese: translation compresses, src_max is the limit
        return src_max
    else:
        # Same script family — use the smaller to be safe
        return min(src_max, tgt_max)


def _split_to_fit(text: str, max_chars: int) -> list[str]:
    """
    Break text into pieces ≤ max_chars.
    For space-delimited text (Latin etc.) splits at word boundaries.
    For CJK and other no-space scripts, falls back to character boundaries.
    """
    if len(text) <= max_chars:
        return [text]
    words = text.split()
    if len(words) > 1:
        # Word-boundary split (Latin, etc.)
        pieces, current = [], ""
        for word in words:
            candidate = (current + " " + word).strip()
            if len(candidate) <= max_chars:
                current = candidate
            else:
                if current:
                    pieces.append(current)
                current = word
        if current:
            pieces.append(current)
        return pieces
    # No spaces (CJK) — split at character boundaries
    return [text[i : i + max_chars] for i in range(0, len(text), max_chars)]


def build_chunks(paragraphs: list[str], mode: str, n: int, max_chars: int = 9999) -> list[str]:
    """
    Greedy packing: fill each slide with as much text as possible before
    starting a new one, so content is consolidated rather than fragmented.

    mode='paragraph' → units are whole paragraphs; short ones are merged
                        together until max_chars is reached.
    mode='sentence'  → units are sentences; pack up to n sentences per slide
                        (or fewer if max_chars would be exceeded).
    max_chars        → hard upper bound on characters per slide chunk.
    """
    if mode == "paragraph":
        units     = paragraphs
        max_units = 999_999   # no count cap — only max_chars limits merging
    else:
        units = []
        for para in paragraphs:
            units.extend(split_sentences(para))
        max_units = n         # respect the "sentences per slide" setting

    chunks: list[str] = []
    current      = ""
    current_count = 0

    for unit in units:
        # If the unit itself is longer than the limit, word-split it first
        pieces = _split_to_fit(unit, max_chars)
        for piece in pieces:
            candidate = (current + " " + piece).strip()
            if current and (len(candidate) > max_chars or current_count >= max_units):
                chunks.append(current)
                current       = piece
                current_count = 1
            else:
                current       = candidate
                current_count += 1

    if current:
        chunks.append(current)

    return chunks

def chars_per_line_for_font(font_pt: int, lang_code: str = "en") -> int:
    """Estimated characters that fit on one line at the given font size."""
    usable_w_pt = (13.33 - 0.9) * 72
    factor      = _char_width_factor(lang_code)
    return max(5, int(usable_w_pt / (font_pt * factor)))


def fix_widow(text: str, chars_per_line: int) -> str:
    """
    Simulate word-wrap. If the last line would contain only one word (widow),
    steal the last word of the previous line and pull it down, then return the
    text with explicit newlines so python-pptx honours the adjusted breaks.
    """
    words = text.split()
    if len(words) < 3:
        return text

    # Simulate word-wrap into lines
    lines: list[list[str]] = [[]]
    for word in words:
        probe = " ".join(lines[-1] + [word])
        if len(probe) <= chars_per_line:
            lines[-1].append(word)
        else:
            lines.append([word])

    if len(lines) < 2 or len(lines[-1]) != 1:
        return text  # no widow — leave untouched (no \n injected)

    if len(lines[-2]) <= 1:
        return text  # previous line too short to donate a word

    # Steal the last word from the previous line
    stolen = lines[-2].pop()
    lines[-1].insert(0, stolen)

    return "\n".join(" ".join(line) for line in lines)


def fix_orphan_chunks(chunks: list[str], max_chars: int, min_words: int = 3) -> list[str]:
    """
    Merge any chunk that is an orphan into the previous chunk when possible.
    For space-delimited text: orphan = fewer than min_words words.
    For CJK (no spaces): orphan = fewer than min_words * 4 characters.
    """
    if len(chunks) <= 1:
        return chunks

    def _is_orphan(text: str) -> bool:
        if " " in text:
            return len(text.split()) < min_words
        return len(text) < min_words * 4   # ~3 short CJK words

    result = list(chunks)
    i = len(result) - 1
    while i > 0:
        if _is_orphan(result[i]):
            merged = (result[i - 1] + " " + result[i]).strip()
            if len(merged) <= max_chars:
                result[i - 1] = merged
                result.pop(i)
            else:
                # Redistribute evenly at the character midpoint
                combined = result[i - 1] + result[i]
                mid = len(combined) // 2
                result[i - 1] = combined[:mid].rstrip()
                result[i]     = combined[mid:].lstrip()
        i -= 1

    return result


# ── Translation ───────────────────────────────────────────────────────────────

def translate_chunks_openai(
    chunks: list[str],
    src_name: str,
    tgt_name: str,
    api_key: str,
    progress_bar,
) -> list[str]:
    client = openai.OpenAI(api_key=api_key)
    results = []
    total = len(chunks)
    system_prompt = (
        f"You are a professional translator. "
        f"Translate the user's text from {src_name} to {tgt_name}. "
        f"Return only the translated text, with no explanations or extra commentary."
    )
    for idx, text in enumerate(chunks):
        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": text},
                ],
                temperature=0.2,
            )
            translated = response.choices[0].message.content.strip() or text
        except Exception as e:
            translated = text  # fall back to original on any error
        results.append(translated)
        progress_bar.progress((idx + 1) / total)
    return results


def translate_chunks_google(
    chunks: list[str],
    src_code: str,
    tgt_code: str,
    progress_bar,
) -> list[str]:
    """Translate using Google Translate (free, no API key required)."""
    results = []
    total = len(chunks)
    # Google Translate uses 'zh-CN' / 'zh-TW'; deep-translator accepts them directly
    for idx, text in enumerate(chunks):
        try:
            translated = GoogleTranslator(source=src_code, target=tgt_code).translate(text) or text
        except Exception:
            translated = text  # fall back to original on any error
        results.append(translated)
        progress_bar.progress((idx + 1) / total)
    return results

# ── PowerPoint builder ────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, color, font_pt, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    # Remove default internal padding so every pixel of the box is usable
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    # Safety net: if text still overflows, PowerPoint shrinks the font to fit
    # rather than letting it bleed into the adjacent language box
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_pt)
        run.font.color.rgb = color
        run.font.bold = bold


def build_pptx(
    chunks: list[str],
    translations: list[str],
    orig_font_pt: int,
    trans_font_pt: int,
    theme_name: str,
    title: str,
    src_lang_code: str = "en",
    tgt_lang_code: str = "zh-CN",
) -> Presentation:
    palette = THEMES[theme_name]
    bg      = palette["bg"]
    orig_c  = palette["orig"]
    trans_c = palette["trans"]
    div_c   = palette["divider"]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # truly blank

    # ── Title slide ──────────────────────────────────────────────────────────
    if title:
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, bg)
        _add_textbox(
            slide,
            Inches(0.5), Inches(2.0),
            Inches(12.33), Inches(3.5),
            title,
            orig_c,
            max(orig_font_pt, 60),
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    # ── Content slides ────────────────────────────────────────────────────────
    # Vertical layout (no labels):
    #   V_PAD  |  orig_box  |  DIV_GAP  |  divider  |  DIV_GAP  |  trans_box  |  V_PAD
    MARGIN  = Inches(0.45)
    TEXT_W  = Inches(13.33) - 2 * MARGIN
    V_PAD   = Inches(0.2)
    DIV_H   = Inches(0.07)
    DIV_GAP = Inches(0.08)

    BOX_H   = (Inches(7.5) - 2 * V_PAD - DIV_H - 2 * DIV_GAP) / 2

    orig_top  = V_PAD
    div_top   = orig_top  + BOX_H + DIV_GAP
    trans_top = div_top   + DIV_H + DIV_GAP

    cpl_src = chars_per_line_for_font(orig_font_pt, src_lang_code)
    cpl_tgt = chars_per_line_for_font(trans_font_pt, tgt_lang_code)

    for orig_text, trans_text in zip(chunks, translations):
        orig_text  = fix_widow(orig_text,  cpl_src)
        trans_text = fix_widow(trans_text, cpl_tgt)
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, bg)

        # Original text (top half)
        _add_textbox(
            slide, MARGIN, orig_top, TEXT_W, BOX_H,
            orig_text, orig_c, orig_font_pt,
        )

        # Divider bar
        divider = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            MARGIN, div_top, TEXT_W, DIV_H,
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = div_c
        divider.line.fill.background()

        # Translated text (bottom half)
        _add_textbox(
            slide, MARGIN, trans_top, TEXT_W, BOX_H,
            trans_text, trans_c, trans_font_pt,
        )

    return prs

# ── Streamlit UI ──────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="Bilingual Sermon Slides",
    layout="wide",
)

st.title("Bilingual Sermon Slide Generator")
st.caption(
    "Upload a Word sermon document and generate bilingual PowerPoint slides "
    "your audience can follow while the sermon is preached."
)

# ── Sidebar settings ──────────────────────────────────────────────────────────
with st.sidebar:
    st.header("Settings")

    translation_engine = st.radio(
        "Translation engine",
        ["Google Translate (free)", "OpenAI (API key required)"],
        help="Google Translate is free and requires no key. OpenAI uses GPT-4o-mini for higher quality.",
    )

    openai_api_key = ""
    if translation_engine == "OpenAI (API key required)":
        openai_api_key = st.text_input(
            "OpenAI API Key",
            type="password",
            placeholder="sk-…",
            help="Your OpenAI API key. Never stored — used only for this session.",
        )

    st.divider()

    src_lang = st.selectbox(
        "Sermon language (source)",
        list(LANGUAGES.keys()),
        index=0,
    )
    tgt_lang = st.selectbox(
        "Audience language (target)",
        list(LANGUAGES.keys()),
        index=1,
    )

    st.divider()

    split_mode = st.radio("Split slides by", ["Sentences", "Paragraphs"])
    sentences_per_slide = 2
    if split_mode == "Sentences":
        sentences_per_slide = st.slider(
            "Sentences per slide", min_value=1, max_value=4, value=2
        )

    st.divider()

    font_mode = st.radio(
        "Font size mode",
        ["Auto (optimal)", "Manual"],
        help="Auto picks the largest font that fits each half-slide. Manual lets you set sizes with sliders.",
    )
    if font_mode == "Manual":
        orig_font_size  = st.slider(f"{src_lang} (original) font size (pt)",   min_value=28, max_value=80, value=54, key="orig_font_size")
        trans_font_size = st.slider(f"{tgt_lang} (translated) font size (pt)", min_value=28, max_value=80, value=54, key="trans_font_size")
    else:
        orig_font_size  = 54  # placeholder — overridden after translation
        trans_font_size = 54
        st.caption("Optimal sizes will be calculated after translation.")
    theme = st.selectbox("Slide theme", list(THEMES.keys()))
    title_text = st.text_input("Presentation title (optional)")

# ── Main area ─────────────────────────────────────────────────────────────────
uploaded_file = st.file_uploader(
    "Upload sermon document (.doc / .docx)",
    type=["doc", "docx"],
    help="Microsoft Word .doc or .docx format. For .doc files, LibreOffice must be installed.",
)

if uploaded_file:
    st.success(f"Loaded: **{uploaded_file.name}**")

    if st.button("Generate Bilingual Slides", type="primary"):

        with st.spinner("Reading document…"):
            try:
                paragraphs = extract_paragraphs(uploaded_file)
            except RuntimeError as e:
                st.error(str(e))
                st.stop()

        if not paragraphs:
            st.error("No text found in the document. Please check the file.")
            st.stop()

        st.info(f"{len(paragraphs)} paragraph(s) read from document.")

        mode_key  = "paragraph" if split_mode == "Paragraphs" else "sentence"
        src_code  = LANGUAGES[src_lang]
        tgt_code  = LANGUAGES[tgt_lang]
        # In auto mode pack text using minimum font so chunks are as dense as possible
        pack_orig  = orig_font_size  if font_mode == "Manual" else 28
        pack_trans = trans_font_size if font_mode == "Manual" else 28
        max_chars = _src_max_chars(pack_orig, pack_trans, src_code, tgt_code)
        chunks = build_chunks(paragraphs, mode_key, sentences_per_slide, max_chars)
        chunks = fix_orphan_chunks(chunks, max_chars)
        st.info(f"{len(chunks)} slide(s) will be created.")

        if translation_engine == "OpenAI (API key required)" and not openai_api_key:
            st.error("Please enter your OpenAI API key in the sidebar before generating slides.")
            st.stop()

        st.write("**Translating…** (this may take a moment for long sermons)")
        progress = st.progress(0)
        if translation_engine == "OpenAI (API key required)":
            translations = translate_chunks_openai(
                chunks,
                src_lang,
                tgt_lang,
                openai_api_key,
                progress,
            )
        else:
            translations = translate_chunks_google(
                chunks,
                src_code,
                tgt_code,
                progress,
            )

        opt_orig  = optimal_font_size(chunks,       src_code)
        opt_trans = optimal_font_size(translations, tgt_code)

        if font_mode == "Auto (optimal)":
            final_orig_font  = opt_orig
            final_trans_font = opt_trans
            st.success(
                f"Auto font sizes — "
                f"**{src_lang}**: {final_orig_font} pt | "
                f"**{tgt_lang}**: {final_trans_font} pt"
            )
        else:
            final_orig_font  = orig_font_size
            final_trans_font = trans_font_size
            st.info(
                f"Optimal font sizes for full slide usage — "
                f"**{src_lang}**: {opt_orig} pt | "
                f"**{tgt_lang}**: {opt_trans} pt"
            )

        with st.spinner("Building PowerPoint…"):
            prs = build_pptx(
                chunks,
                translations,
                orig_font_pt=final_orig_font,
                trans_font_pt=final_trans_font,
                theme_name=theme,
                title=title_text,
                src_lang_code=src_code,
                tgt_lang_code=tgt_code,
            )
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)

        st.success(f"Done! {len(chunks)} slides generated.")
        st.download_button(
            label="Download PowerPoint (.pptx)",
            data=buf,
            file_name="sermon_bilingual.pptx",
            mime=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        )
